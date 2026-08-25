# This is copied from Magentic-one's great repo: https://github.com/microsoft/autogen/blob/v0.4.4/python/packages/autogen-magentic-one/src/autogen_magentic_one/markdown_browser/mdconvert.py
# Thanks to Microsoft researchers for open-sourcing this!
# type: ignore
import base64
import copy
import html
import json
import mimetypes
import os
import re
import shutil
import subprocess
import sys
import tempfile
import traceback
import zipfile
import uuid
from pathlib import Path
from openpyxl import load_workbook
from typing import Any, Dict, List, Optional, Union
from urllib.parse import parse_qs, quote, unquote, urlparse, urlunparse
from openai import OpenAI

import mammoth
import markdownify
import pandas as pd
import pdfminer
import pdfminer.high_level
import pptx

# File-format detection
import puremagic
import pydub
import requests
import speech_recognition as sr
from bs4 import BeautifulSoup
from youtube_transcript_api import YouTubeTranscriptApi
from youtube_transcript_api.formatters import SRTFormatter
from supadata import Supadata, SupadataError


class _CustomMarkdownify(markdownify.MarkdownConverter):
    """
    A custom version of markdownify's MarkdownConverter. Changes include:

    - Altering the default heading style to use '#', '##', etc.
    - Removing javascript hyperlinks.
    - Truncating images with large data:uri sources.
    - Ensuring URIs are properly escaped, and do not conflict with Markdown syntax
    """

    def __init__(self, **options: Any):
        options["heading_style"] = options.get("heading_style", markdownify.ATX)
        # Explicitly cast options to the expected type if necessary
        super().__init__(**options)

    def convert_hn(self, n: int, el: Any, text: str, convert_as_inline: bool) -> str:
        """Same as usual, but be sure to start with a new line"""
        if not convert_as_inline:
            if not re.search(r"^\n", text):
                return "\n" + super().convert_hn(n, el, text, convert_as_inline)  # type: ignore

        return super().convert_hn(n, el, text, convert_as_inline)  # type: ignore

    def convert_a(self, el: Any, text: str, convert_as_inline: bool):
        """Same as usual converter, but removes Javascript links and escapes URIs."""
        prefix, suffix, text = markdownify.chomp(text)  # type: ignore
        if not text:
            return ""
        href = el.get("href")
        title = el.get("title")

        # Escape URIs and skip non-http or file schemes
        if href:
            try:
                parsed_url = urlparse(href)  # type: ignore
                if parsed_url.scheme and parsed_url.scheme.lower() not in ["http", "https", "file"]:  # type: ignore
                    return "%s%s%s" % (prefix, text, suffix)
                href = urlunparse(parsed_url._replace(path=quote(unquote(parsed_url.path))))  # type: ignore
            except ValueError:  # It's not clear if this ever gets thrown
                return "%s%s%s" % (prefix, text, suffix)

        # For the replacement see #29: text nodes underscores are escaped
        if (
            self.options["autolinks"]
            and text.replace(r"\_", "_") == href
            and not title
            and not self.options["default_title"]
        ):
            # Shortcut syntax
            return "<%s>" % href
        if self.options["default_title"] and not title:
            title = href
        title_part = ' "%s"' % title.replace('"', r"\"") if title else ""
        return "%s[%s](%s%s)%s" % (prefix, text, href, title_part, suffix) if href else text

    def convert_img(self, el: Any, text: str, convert_as_inline: bool) -> str:
        """Same as usual converter, but removes data URIs"""

        alt = el.attrs.get("alt", None) or ""
        src = el.attrs.get("src", None) or ""
        title = el.attrs.get("title", None) or ""
        title_part = ' "%s"' % title.replace('"', r"\"") if title else ""
        if convert_as_inline and el.parent.name not in self.options["keep_inline_images_in"]:
            return alt

        # Remove dataURIs
        if src.startswith("data:"):
            src = src.split(",")[0] + "..."

        return "![%s](%s%s)" % (alt, src, title_part)

    def convert_soup(self, soup: Any) -> str:
        return super().convert_soup(soup)  # type: ignore


class DocumentConverterResult:
    """The result of converting a document to text."""

    def __init__(self, title: Union[str, None] = None, text_content: str = ""):
        self.title: Union[str, None] = title
        self.text_content: str = text_content


class DocumentConverter:
    """Abstract superclass of all DocumentConverters."""

    def convert(self, local_path: str, **kwargs: Any) -> Union[None, DocumentConverterResult]:
        raise NotImplementedError()


class PlainTextConverter(DocumentConverter):
    """Anything with content type text/plain"""

    def convert(self, local_path: str, **kwargs: Any) -> Union[None, DocumentConverterResult]:
        # Guess the content type from any file extension that might be around
        content_type, _ = mimetypes.guess_type("__placeholder" + kwargs.get("file_extension", ""))

        # Only accept text files
        if content_type is None:
            return None
        # elif "text/" not in content_type.lower():
        #     return None

        text_content = ""
        with open(local_path, "rt", encoding="utf-8") as fh:
            text_content = fh.read()
        return DocumentConverterResult(
            title=None,
            text_content=text_content,
        )


class HtmlConverter(DocumentConverter):
    """Anything with content type text/html"""

    def convert(self, local_path: str, **kwargs: Any) -> Union[None, DocumentConverterResult]:
        # Bail if not html
        extension = kwargs.get("file_extension", "")
        if extension.lower() not in [".html", ".htm"]:
            return None

        result = None
        with open(local_path, "rt", encoding="utf-8") as fh:
            result = self._convert(fh.read())

        return result

    def _convert(self, html_content: str) -> Union[None, DocumentConverterResult]:
        """Helper function that converts and HTML string."""

        # Parse the string
        soup = BeautifulSoup(html_content, "html.parser")

        # Remove javascript and style blocks
        for script in soup(["script", "style"]):
            script.extract()

        # Print only the main content
        body_elm = soup.find("body")
        webpage_text = ""
        if body_elm:
            webpage_text = _CustomMarkdownify().convert_soup(body_elm)
        else:
            webpage_text = _CustomMarkdownify().convert_soup(soup)

        assert isinstance(webpage_text, str)

        return DocumentConverterResult(
            title=None if soup.title is None else soup.title.string, text_content=webpage_text
        )


class WikipediaConverter(DocumentConverter):
    """Handle Wikipedia pages separately, focusing only on the main document content."""

    def convert(self, local_path: str, **kwargs: Any) -> Union[None, DocumentConverterResult]:
        # Bail if not Wikipedia
        extension = kwargs.get("file_extension", "")
        if extension.lower() not in [".html", ".htm"]:
            return None
        url = kwargs.get("url", "")
        if not re.search(r"^https?:\/\/[a-zA-Z]{2,3}\.wikipedia.org\/", url):
            return None

        # Parse the file
        soup = None
        with open(local_path, "rt", encoding="utf-8") as fh:
            soup = BeautifulSoup(fh.read(), "html.parser")

        # Remove javascript and style blocks
        for script in soup(["script", "style"]):
            script.extract()

        # Print only the main content
        body_elm = soup.find("div", {"id": "mw-content-text"})
        title_elm = soup.find("span", {"class": "mw-page-title-main"})

        webpage_text = ""
        main_title = None if soup.title is None else soup.title.string

        if body_elm:
            # What's the title
            if title_elm and len(title_elm) > 0:
                main_title = title_elm.string  # type: ignore
                assert isinstance(main_title, str)

            # Convert the page
            webpage_text = f"# {main_title}\n\n" + _CustomMarkdownify().convert_soup(body_elm)
        else:
            webpage_text = _CustomMarkdownify().convert_soup(soup)

        return DocumentConverterResult(
            title=main_title,
            text_content=webpage_text,
        )


class YouTubeConverter(DocumentConverter):
    """Handle YouTube specially, focusing on the video title, description, transcript, and comments."""

    def convert(self, local_path: str, **kwargs: Any) -> Union[None, DocumentConverterResult]:
        # Bail if not YouTube
        extension = kwargs.get("file_extension", "")
        if extension.lower() not in [".html", ".htm"]:
            return None
        url = kwargs.get("url", "")
        if not url.startswith("https://www.youtube.com/watch?"):
            return None
            
        # Check if API key is available
        api_key = os.environ.get("SUPADATA_API_KEY")
        if api_key:
            return self._convert_using_api(url, local_path)
        else:
            return self._convert_using_local(url, local_path)
            
    def _fetch_comments(self, video_id: str) -> str:
        """Fetch top 100 comments from a YouTube video."""
        try:
            api_key = os.environ.get("YOUTUBE_API_KEY")
            if not api_key:
                return ""
                
            url = f"https://www.googleapis.com/youtube/v3/commentThreads"
            params = {
                "part": "snippet",
                "videoId": video_id,
                "key": api_key,
                "maxResults": 100
            }
            
            response = requests.get(url, params=params)
            if response.status_code != 200:
                print(f"Failed to fetch comments: {response.status_code}")
                return ""
                
            comments_data = response.json()
            comments_text = ""
            
            for item in comments_data.get("items", []):
                comment = item["snippet"]["topLevelComment"]["snippet"]["textOriginal"]
                comments_text += f"{comment}\n\n"
                
            return comments_text
        except Exception as e:
            print(f"Error fetching comments: {e}")
            return ""
    
    def _convert_using_api(self, url: str, local_path: str) -> Union[None, DocumentConverterResult]:
        """Use Supadata API to get YouTube data."""
        try:
            api_key = os.environ.get("SUPADATA_API_KEY")
            if not api_key:
                return None
                
            supadata_client = Supadata(api_key=api_key)
            
            transcript_text = supadata_client.transcript(
                url=url,
                lang="en",
                text=True,
                mode="auto"
            )
            
            parsed_url = urlparse(url)
            params = parse_qs(parsed_url.query)
            if "v" not in params:
                return None
                 
            video_id = str(params["v"][0])
            webpage_text = "# YouTube\n"
            title = f"YouTube Video: {video_id}"
            webpage_text += f"\n## {title}\n"
            
            if transcript_text:
                webpage_text += f"\n### Transcript\n{transcript_text}\n"
            
            # Fetch and add comments
            comments_text = self._fetch_comments(video_id)
            if comments_text:
                webpage_text += f"\n### Top Comments\n{comments_text}\n"
                
            return DocumentConverterResult(
                title=title,
                text_content=webpage_text,
            )
            
        except Exception as e:
            print(f"API request failed: {e}")
            print("Falling back to local conversion method...")
            return self._convert_using_local(url, local_path)
            
    def _convert_using_local(self, url: str, local_path: str) -> Union[None, DocumentConverterResult]:
        """Use the original local conversion method."""
        # Parse the file
        soup = None
        with open(local_path, "rt", encoding="utf-8") as fh:
            soup = BeautifulSoup(fh.read(), "html.parser")

        # Read the meta tags
        assert soup.title is not None and soup.title.string is not None
        metadata: Dict[str, str] = {"title": soup.title.string}
        for meta in soup(["meta"]):
            for a in meta.attrs:
                if a in ["itemprop", "property", "name"]:
                    metadata[meta[a]] = meta.get("content", "")
                    break

        # We can also try to read the full description. This is more prone to breaking, since it reaches into the page implementation
        try:
            for script in soup(["script"]):
                content = script.text
                if "ytInitialData" in content:
                    lines = re.split(r"\r?\n", content)
                    obj_start = lines[0].find("{")
                    obj_end = lines[0].rfind("}")
                    if obj_start >= 0 and obj_end >= 0:
                        data = json.loads(lines[0][obj_start : obj_end + 1])
                        attrdesc = self._findKey(data, "attributedDescriptionBodyText")  # type: ignore
                        if attrdesc:
                            metadata["description"] = str(attrdesc["content"])
                    break
        except Exception:
            pass

        # Start preparing the page
        webpage_text = "# YouTube\n"

        title = self._get(metadata, ["title", "og:title", "name"])  # type: ignore
        assert isinstance(title, str)

        if title:
            webpage_text += f"\n## {title}\n"

        stats = ""
        views = self._get(metadata, ["interactionCount"])  # type: ignore
        if views:
            stats += f"- **Views:** {views}\n"

        keywords = self._get(metadata, ["keywords"])  # type: ignore
        if keywords:
            stats += f"- **Keywords:** {keywords}\n"

        runtime = self._get(metadata, ["duration"])  # type: ignore
        if runtime:
            stats += f"- **Runtime:** {runtime}\n"

        if len(stats) > 0:
            webpage_text += f"\n### Video Metadata\n{stats}\n"

        description = self._get(metadata, ["description", "og:description"])  # type: ignore
        if description:
            webpage_text += f"\n### Description\n{description}\n"

        transcript_text = ""
        parsed_url = urlparse(url)  # type: ignore
        params = parse_qs(parsed_url.query)  # type: ignore
        if "v" in params:
            assert isinstance(params["v"][0], str)
            video_id = str(params["v"][0])
            try:
                # Must be a single transcript.
                transcript = YouTubeTranscriptApi.get_transcript(video_id)  # type: ignore
                # transcript_text = " ".join([part["text"] for part in transcript])  # type: ignore
                # Alternative formatting:
                transcript_text = SRTFormatter().format_transcript(transcript)
            except Exception:
                pass
        if transcript_text:
            webpage_text += f"\n### Transcript\n{transcript_text}\n"
            
        # Fetch and add comments if video_id is available
        if "v" in params:
            video_id = str(params["v"][0])
            comments_text = self._fetch_comments(video_id)
            if comments_text:
                webpage_text += f"\n### Top Comments\n{comments_text}\n"

        title = title if title else soup.title.string
        assert isinstance(title, str)

        return DocumentConverterResult(
            title=title,
            text_content=webpage_text,
        )

    def _get(self, metadata: Dict[str, str], keys: List[str], default: Union[str, None] = None) -> Union[str, None]:
        for k in keys:
            if k in metadata:
                return metadata[k]
        return default

    def _findKey(self, json: Any, key: str) -> Union[str, None]:  # TODO: Fix json type
        if isinstance(json, list):
            for elm in json:
                ret = self._findKey(elm, key)
                if ret is not None:
                    return ret
        elif isinstance(json, dict):
            for k in json:
                if k == key:
                    return json[k]
                else:
                    ret = self._findKey(json[k], key)
                    if ret is not None:
                        return ret
        return None


class PdfConverter(DocumentConverter):
    """
    Converts PDFs to Markdown. Most style information is ignored, so the results are essentially plain-text.
    """

    def convert(self, local_path, **kwargs) -> Union[None, DocumentConverterResult]:
        # Bail if not a PDF
        extension = kwargs.get("file_extension", "")
        if extension.lower() != ".pdf":
            return None

        return DocumentConverterResult(
            title=None,
            text_content=pdfminer.high_level.extract_text(local_path),
        )


class DocxConverter(HtmlConverter):
    """
    Converts DOCX files to Markdown. Style information (e.g.m headings) and tables are preserved where possible.
    """

    def convert(self, local_path, **kwargs) -> Union[None, DocumentConverterResult]:
        # Bail if not a DOCX
        extension = kwargs.get("file_extension", "")
        if extension.lower() != ".docx":
            return None

        result = None
        with open(local_path, "rb") as docx_file:
            result = mammoth.convert_to_html(docx_file)
            html_content = result.value
            result = self._convert(html_content)

        return result

from openpyxl.utils.exceptions import InvalidFileException
class XlsxConverter(HtmlConverter):
    """
    Converts XLSX files to Markdown, with each sheet presented as a separate Markdown table.
    Marking cells that have background color(only if there are 2 or more distinct colors in the sheet.
    """
    NONE_COLOR_TAG = "none"  

    def convert(self, local_path, **kwargs) -> Union[None, "DocumentConverterResult"]:
        extension = (kwargs.get("file_extension") or "").lower()
        if extension not in (".xlsx", ".xls"):
            return None

        if extension == ".xls":
            try:
                sheets = pd.read_excel(local_path, sheet_name=None)
            except Exception:
                return None
            md_content = ""
            for s in sheets:
                md_content += f"## {s}\n"
                html_content = sheets[s].to_html(index=False)
                md_content += self._convert(html_content).text_content.strip() + "\n\n"

            return DocumentConverterResult(
                title=None,
                text_content=md_content.strip(),
            )

        wb = load_workbook(local_path, data_only=False, read_only=False)
        md_parts: List[str] = []

        for sheet_name in wb.sheetnames:
            ws = wb[sheet_name]

            max_row, max_col, color_categories = self._scan_extent_and_colors(ws)
            if max_row == 0 or max_col == 0:
                continue  
            show_colors = len(color_categories) >= 2  
            rows: List[List[str]] = []
            for r in range(1, max_row + 1):
                row_cells = []
                for c in range(1, max_col + 1):
                    cell = ws.cell(row=r, column=c)
                    row_cells.append(self._render_cell(cell, show_colors))
                rows.append(row_cells)
            md_table = self._rows_to_markdown(rows)
            md_parts.append(f"## {sheet_name}\n{md_table}")

        md_content = "\n\n".join(md_parts).strip()
        return DocumentConverterResult(title=None, text_content=md_content)

    def _scan_extent_and_colors(self, ws):
        max_row = 0
        max_col = 0
        color_categories = set()

        for r in range(1, ws.max_row + 1):
            row_has_data = False
            last_used_col = 0
            for c in range(1, ws.max_column + 1):
                cell = ws.cell(row=r, column=c)
                value_blank = self._is_blank_value(cell.value)
                color_hex = self._extract_bg_color(cell)
                cat = color_hex if color_hex else self.NONE_COLOR_TAG
                color_categories.add(cat)
                if (not value_blank) or color_hex:
                    row_has_data = True
                    last_used_col = c
            if row_has_data:
                max_row = r
                if last_used_col > max_col:
                    max_col = last_used_col

        return max_row, max_col, color_categories

    def _render_cell(self, cell, show_colors: bool) -> str:
        val = cell.value
        txt = "" if self._is_blank_value(val) else str(val)
        color_hex = self._extract_bg_color(cell)

        if show_colors:
            tag = color_hex if color_hex else self.NONE_COLOR_TAG
            return f"{txt} [{tag}]" if txt else f"[{tag}]"
        else:
            return txt

    def _is_blank_value(self, value) -> bool:
        if value is None or value == "":
            return True
        try:
            return pd.isna(value)
        except Exception:
            return False

    def _extract_bg_color(self, cell):
        """
        判断是否“有背景填充”：
        - fill.fill_type (或 patternType) 必须非 None/'none'
        - 尝试读取 start_color.rgb (若无则 fgColor.rgb)
        - 返回 '#RRGGBB'；解析失败返回 None
        *白色视作None*
        """
        fill = getattr(cell, "fill", None)
        if not fill:
            return None

        ftype = getattr(fill, "fill_type", None) or getattr(fill, "patternType", None)
        if not ftype or str(ftype).lower() == "none":
            return None

        color_obj = getattr(fill, "start_color", None) or getattr(fill, "fgColor", None)
        if not color_obj:
            return None

        rgb_val = getattr(color_obj, "rgb", None)
        if not rgb_val:
            return None

        hex_str = str(rgb_val)

        if len(hex_str) == 8:     
            return "#" + hex_str[2:] if hex_str[2:] != "FFFFFF" else None
        elif len(hex_str) == 6:   
            return "#" + hex_str if hex_str != "FFFFFF" else None
        else:
            return None


    def _rows_to_markdown(self, rows):
        if not rows:
            return ""
        header = rows[0]
        lines = [
            "| " + " | ".join(header) + " |",
            "| " + " | ".join(["---"] * len(header)) + " |",
        ]
        for row in rows[1:]:
            lines.append("| " + " | ".join(row) + " |")
        return "\n".join(lines)


class PptxConverter(HtmlConverter):
    """
    Converts PPTX files to Markdown. Supports heading, tables and images with alt text.
    """

    def __init__(self, extract_dir: str = "downloads/ppt_pics") -> None:
        self.extract_dir = extract_dir
        os.makedirs(self.extract_dir, exist_ok=True)
        super().__init__()

    def convert(self, local_path, **kwargs: Any) -> Union[None, "DocumentConverterResult"]:
        extension = kwargs.get("file_extension", "")
        if extension.lower() != ".pptx":
            return None

        presentation = pptx.Presentation(local_path)
        md_chunks: List[str] = []
        extracted_files: List[str] = []

        for slide_idx, slide in enumerate(presentation.slides, start=1):
            md_chunks.append(f"\n\n<!-- Slide number: {slide_idx} -->\n")
            title_shape = slide.shapes.title if hasattr(slide.shapes, "title") else None
            for shape_idx, shape in enumerate(slide.shapes, start=1):
                if self._is_picture(shape):
                    alt_text = self._get_alt_text(shape)
                    img_path = self._extract_image(shape, slide_idx, shape_idx)
                    if img_path:
                        extracted_files.append(img_path)
                        md_chunks.append(f"\n![{alt_text}]({img_path})\n")
                    else:
                        filename = re.sub(r"\W", "", shape.name) + ".jpg"
                        md_chunks.append(f"\n![{alt_text}]({filename})\n")

                if self._is_table(shape):
                    html_table = ["<html><body><table>"]
                    first_row = True
                    for row in shape.table.rows:
                        html_table.append("<tr>")
                        for cell in row.cells:
                            cell_txt = html.escape(cell.text)
                            if first_row:
                                html_table.append(f"<th>{cell_txt}</th>")
                            else:
                                html_table.append(f"<td>{cell_txt}</td>")
                        html_table.append("</tr>")
                        first_row = False
                    html_table.append("</table></body></html>")

                    html_blob = "".join(html_table)
                    md_table = self._convert(html_blob).text_content.strip()
                    md_chunks.append("\n" + md_table + "\n")

                elif getattr(shape, "has_text_frame", False):
                    text_val = shape.text or ""
                    if title_shape is not None and shape == title_shape:
                        md_chunks.append("# " + text_val.lstrip() + "\n")
                    else:
                        md_chunks.append(text_val + "\n")

            if slide.has_notes_slide:
                notes_frame = slide.notes_slide.notes_text_frame
                if notes_frame is not None:
                    md_chunks.append("\n\n### Notes:\n")
                    md_chunks.append(notes_frame.text)

        md_content = "".join(md_chunks).strip()

        return DocumentConverterResult(
            title=None,
            text_content=md_content,
        )

    def _is_picture(self, shape) -> bool:
        mso = pptx.enum.shapes.MSO_SHAPE_TYPE
        if shape.shape_type == mso.PICTURE:
            return True
        if shape.shape_type == mso.PLACEHOLDER and hasattr(shape, "image"):
            return True
        return False

    def _is_table(self, shape) -> bool:
        mso = pptx.enum.shapes.MSO_SHAPE_TYPE
        return shape.shape_type == mso.TABLE

    def _get_alt_text(self, shape) -> str:
        alt_text = ""
        try:
            alt_text = shape._element._nvXxPr.cNvPr.attrib.get("descr", "")  
        except Exception:
            pass
        if not alt_text:
            alt_text = getattr(shape, "name", "image") or "image"
        return alt_text

    def _extract_image(self, shape, slide_idx: int, shape_idx: int) -> Optional[str]:
        """Extract the binary for a picture *if available*; return relative path or None.

        We try to preserve the native extension from ``shape.image.ext`` when present.
        Filenames are normalized to ``slide{n}_shape{m}.{ext}`` to avoid collisions and
        to keep paths deterministic (important for caching / testing). All files land in
        ``self.extract_dir``.
        """
        image = getattr(shape, "image", None)
        if image is None:
            return None

        ext = getattr(image, "ext", None)
        if ext:
            ext = ext.lower().lstrip('.')
        else:
            ctype = getattr(image, "content_type", "").lower()
            if "/" in ctype:
                ext = ctype.split("/")[-1]
            else:
                ext = "img" 
        if ext == "jpeg":
            ext = "jpg"
        fname = f"slide{slide_idx}_shape{shape_idx}_{uuid.uuid4().hex[:8]}.{ext}"
        out_path = Path(self.extract_dir) / fname
        try:
            with open(out_path, "wb") as f:
                f.write(image.blob)
        except Exception:
            return None
        return out_path.as_posix()


class MediaConverter(DocumentConverter):
    """
    Abstract class for multi-modal media (e.g., images and audio)
    """

    def _get_metadata(self, local_path):
        exiftool = shutil.which("exiftool")
        if not exiftool:
            return None
        else:
            try:
                result = subprocess.run([exiftool, "-json", local_path], capture_output=True, text=True).stdout
                return json.loads(result)[0]
            except Exception:
                return None


class WavConverter(MediaConverter):
    """
    Converts WAV files to markdown via extraction of metadata (if `exiftool` is installed), and speech transcription (if `speech_recognition` is installed).
    """

    def convert(self, local_path, **kwargs) -> Union[None, DocumentConverterResult]:
        extension = kwargs.get("file_extension", "")
        if extension.lower() != ".wav":
            return None

        md_content = ""

        metadata = self._get_metadata(local_path)
        if metadata:
            for f in [
                "Title",
                "Artist",
                "Author",
                "Band",
                "Album",
                "Genre",
                "Track",
                "DateTimeOriginal",
                "CreateDate",
                "Duration",
            ]:
                if f in metadata:
                    md_content += f"{f}: {metadata[f]}\n"
        try:
            transcript = self._transcribe_audio(local_path)
            md_content += "\n\n### Audio Transcript:\n" + ("[No speech detected]" if transcript == "" else transcript)
        except Exception:
            md_content += "\n\n### Audio Transcript:\nError. Could not transcribe this audio."

        return DocumentConverterResult(
            title=None,
            text_content=md_content.strip(),
        )

    def _transcribe_audio(self, local_path) -> str:
        recognizer = sr.Recognizer()
        with sr.AudioFile(local_path) as source:
            audio = recognizer.record(source)
            return recognizer.recognize_google(audio).strip()


class Mp3Converter(WavConverter):
    """
    Converts MP3 and M4A files to markdown via extraction of metadata (if `exiftool` is installed), and speech transcription (if `speech_recognition` AND `pydub` are installed).
    """

    def __init__(self):
        self.api_key = os.getenv("PLUGIN_API_KEY")
        self.base_url = os.getenv("PLUGIN_BASE_URL")

    # This is copy from Agent-KB repo
    def transcribe_audio_with_api(self, file_path: str) -> str:
        """Transcribe audio using OpenAI Whisper API."""
        client = OpenAI(api_key=self.api_key, base_url=self.base_url)
        try:
            with open(file_path, "rb") as audio_file:
                transcription = client.audio.transcriptions.create(
                    model="whisper-1",
                    file=audio_file
                )
            return transcription.text
        except Exception as e:
            raise RuntimeError(f"Speech recognition via API failed: {str(e)}") from e

    def convert(self, local_path, **kwargs) -> Union[None, DocumentConverterResult]:
        extension = kwargs.get("file_extension", "")
        if extension.lower() not in [".mp3", ".m4a"]:
            return None

        md_content = ""

        metadata = self._get_metadata(local_path)
        if metadata:
            for f in [
                "Title", "Artist", "Author", "Band", "Album", "Genre",
                "Track", "DateTimeOriginal", "CreateDate", "Duration",
            ]:
                if f in metadata:
                    md_content += f"{f}: {metadata[f]}\n"

        handle, temp_path = tempfile.mkstemp(suffix=".wav")
        os.close(handle)
        try:
            if extension.lower() == ".mp3":
                sound = pydub.AudioSegment.from_mp3(local_path)
            else:
                sound = pydub.AudioSegment.from_file(local_path, format="m4a")
            sound.export(temp_path, format="wav")

            # Use API if credentials exist, else fallback
            try:
                if self.api_key and self.base_url:
                    transcript = self.transcribe_audio_with_api(temp_path).strip()
                else:
                    print("OpenAI api_key not found, Using fallback transcription method")
                    transcript = super()._transcribe_audio(temp_path).strip()

                md_content += "\n\n### Audio Transcript:\n" + (
                    "[No speech detected]" if transcript == "" else transcript
                )
            except Exception:
                md_content += "\n\n### Audio Transcript:\nError. Could not transcribe this audio."

        finally:
            os.unlink(temp_path)

        return DocumentConverterResult(
            title=None,
            text_content=md_content.strip(),
        )


class ZipConverter(DocumentConverter):
    """
    Extracts ZIP files to a permanent local directory and returns a listing of extracted files.
    """

    def __init__(self, extract_dir: str = "downloads"):
        self.extract_dir = extract_dir
        os.makedirs(self.extract_dir, exist_ok=True)

    def convert(self, local_path: str, **kwargs: Any) -> Union[None, DocumentConverterResult]:
        # Bail if not a ZIP file
        extension = kwargs.get("file_extension", "")
        if extension.lower() != ".zip":
            return None

        # Verify it's actually a ZIP file
        if not zipfile.is_zipfile(local_path):
            return None

        # Extract all files and build list
        extracted_files = []
        with zipfile.ZipFile(local_path, "r") as zip_ref:
            # Extract all files
            zip_ref.extractall(self.extract_dir)
            # Get list of all files
            for file_path in zip_ref.namelist():
                # Skip directories
                if not file_path.endswith("/"):
                    extracted_files.append(self.extract_dir + "/" + file_path)

        # Sort files for consistent output
        extracted_files.sort()

        # Build the markdown content
        md_content = "Downloaded the following files:\n"
        for file in extracted_files:
            md_content += f"* {file}\n"

        return DocumentConverterResult(title="Extracted Files", text_content=md_content.strip())


class ImageConverter(MediaConverter):
    """
    Converts images to markdown via extraction of metadata (if `exiftool` is installed), OCR (if `easyocr` is installed), and description via a multimodal LLM (if an mlm_client is configured).
    """

    def convert(self, local_path, **kwargs) -> Union[None, DocumentConverterResult]:
        # Bail if not a XLSX
        extension = kwargs.get("file_extension", "")
        if extension.lower() not in [".jpg", ".jpeg", ".png"]:
            return None

        md_content = ""

        # Add metadata
        metadata = self._get_metadata(local_path)
        if metadata:
            for f in [
                "ImageSize",
                "Title",
                "Caption",
                "Description",
                "Keywords",
                "Artist",
                "Author",
                "DateTimeOriginal",
                "CreateDate",
                "GPSPosition",
            ]:
                if f in metadata:
                    md_content += f"{f}: {metadata[f]}\n"

        # Try describing the image with GPTV
        mlm_client = kwargs.get("mlm_client")
        mlm_model = kwargs.get("mlm_model")
        if mlm_client is not None and mlm_model is not None:
            md_content += (
                "\n# Description:\n"
                + self._get_mlm_description(
                    local_path, extension, mlm_client, mlm_model, prompt=kwargs.get("mlm_prompt")
                ).strip()
                + "\n"
            )

        return DocumentConverterResult(
            title=None,
            text_content=md_content,
        )

    def _get_mlm_description(self, local_path, extension, client, model, prompt=None):
        if prompt is None or prompt.strip() == "":
            prompt = "Write a detailed caption for this image."

        sys.stderr.write(f"MLM Prompt:\n{prompt}\n")

        data_uri = ""
        with open(local_path, "rb") as image_file:
            content_type, encoding = mimetypes.guess_type("_dummy" + extension)
            if content_type is None:
                content_type = "image/jpeg"
            image_base64 = base64.b64encode(image_file.read()).decode("utf-8")
            data_uri = f"data:{content_type};base64,{image_base64}"

        messages = [
            {
                "role": "user",
                "content": [
                    {"type": "text", "text": prompt},
                    {
                        "type": "image_url",
                        "image_url": {
                            "url": data_uri,
                        },
                    },
                ],
            }
        ]

        response = client.chat.completions.create(model=model, messages=messages)
        return response.choices[0].message.content


class FileConversionException(Exception):
    pass


class UnsupportedFormatException(Exception):
    pass


class MarkdownConverter:
    """(In preview) An extremely simple text-based document reader, suitable for LLM use.
    This reader will convert common file-types or webpages to Markdown."""

    def __init__(
        self,
        requests_session: Optional[requests.Session] = None,
        mlm_client: Optional[Any] = None,
        mlm_model: Optional[Any] = None,
    ):
        if requests_session is None:
            self._requests_session = requests.Session()
        else:
            self._requests_session = requests_session

        self._mlm_client = mlm_client
        self._mlm_model = mlm_model

        self._page_converters: List[DocumentConverter] = []

        # Register converters for successful browsing operations
        # Later registrations are tried first / take higher priority than earlier registrations
        # To this end, the most specific converters should appear below the most generic converters
        self.register_page_converter(PlainTextConverter())
        self.register_page_converter(HtmlConverter())
        self.register_page_converter(WikipediaConverter())
        self.register_page_converter(YouTubeConverter())
        self.register_page_converter(DocxConverter())
        self.register_page_converter(XlsxConverter())
        self.register_page_converter(PptxConverter())
        self.register_page_converter(WavConverter())
        self.register_page_converter(Mp3Converter())
        self.register_page_converter(ImageConverter())
        self.register_page_converter(ZipConverter())
        self.register_page_converter(PdfConverter())

    def convert(
        self, source: Union[str, requests.Response], **kwargs: Any
    ) -> DocumentConverterResult:  # TODO: deal with kwargs
        """
        Args:
            - source: can be a string representing a path or url, or a requests.response object
            - extension: specifies the file extension to use when interpreting the file. If None, infer from source (path, uri, content-type, etc.)
        """

        # Local path or url
        if isinstance(source, str):
            if source.startswith("http://") or source.startswith("https://") or source.startswith("file://"):
                return self.convert_url(source, **kwargs)
            else:
                return self.convert_local(source, **kwargs)
        # Request response
        elif isinstance(source, requests.Response):
            return self.convert_response(source, **kwargs)

    def convert_local(self, path: str, **kwargs: Any) -> DocumentConverterResult:  # TODO: deal with kwargs
        # Prepare a list of extensions to try (in order of priority)
        ext = kwargs.get("file_extension")
        extensions = [ext] if ext is not None else []

        # Get extension alternatives from the path and puremagic
        base, ext = os.path.splitext(path)
        self._append_ext(extensions, ext)
        self._append_ext(extensions, self._guess_ext_magic(path))

        # Convert
        return self._convert(path, extensions, **kwargs)

    # TODO what should stream's type be?
    def convert_stream(self, stream: Any, **kwargs: Any) -> DocumentConverterResult:  # TODO: deal with kwargs
        # Prepare a list of extensions to try (in order of priority)
        ext = kwargs.get("file_extension")
        extensions = [ext] if ext is not None else []

        # Save the file locally to a temporary file. It will be deleted before this method exits
        handle, temp_path = tempfile.mkstemp()
        fh = os.fdopen(handle, "wb")
        result = None
        try:
            # Write to the temporary file
            content = stream.read()
            if isinstance(content, str):
                fh.write(content.encode("utf-8"))
            else:
                fh.write(content)
            fh.close()

            # Use puremagic to check for more extension options
            self._append_ext(extensions, self._guess_ext_magic(temp_path))

            # Convert
            result = self._convert(temp_path, extensions, **kwargs)
        # Clean up
        finally:
            try:
                fh.close()
            except Exception:
                pass
            os.unlink(temp_path)

        return result

    def convert_url(self, url: str, **kwargs: Any) -> DocumentConverterResult:  # TODO: fix kwargs type
        # Send a HTTP request to the URL
        user_agent = "Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36 (KHTML, like Gecko) Chrome/119.0.0.0 Safari/537.36 Edg/119.0.0.0"
        response = self._requests_session.get(url, stream=True, headers={"User-Agent": user_agent})
        response.raise_for_status()
        return self.convert_response(response, **kwargs)

    def convert_response(
        self, response: requests.Response, **kwargs: Any
    ) -> DocumentConverterResult:  # TODO fix kwargs type
        # Prepare a list of extensions to try (in order of priority)
        ext = kwargs.get("file_extension")
        extensions = [ext] if ext is not None else []

        # Guess from the mimetype
        content_type = response.headers.get("content-type", "").split(";")[0]
        self._append_ext(extensions, mimetypes.guess_extension(content_type))

        # Read the content disposition if there is one
        content_disposition = response.headers.get("content-disposition", "")
        m = re.search(r"filename=([^;]+)", content_disposition)
        if m:
            base, ext = os.path.splitext(m.group(1).strip("\"'"))
            self._append_ext(extensions, ext)

        # Read from the extension from the path
        base, ext = os.path.splitext(urlparse(response.url).path)
        self._append_ext(extensions, ext)

        # Save the file locally to a temporary file. It will be deleted before this method exits
        handle, temp_path = tempfile.mkstemp()
        fh = os.fdopen(handle, "wb")
        result = None
        try:
            # Download the file
            for chunk in response.iter_content(chunk_size=512):
                fh.write(chunk)
            fh.close()

            # Use puremagic to check for more extension options
            self._append_ext(extensions, self._guess_ext_magic(temp_path))

            # Convert
            result = self._convert(temp_path, extensions, url=response.url)
        except Exception as e:
            print(f"Error in converting: {e}")

        # Clean up
        finally:
            try:
                fh.close()
            except Exception:
                pass
            os.unlink(temp_path)

        return result

    def _convert(self, local_path: str, extensions: List[Union[str, None]], **kwargs) -> DocumentConverterResult:
        error_trace = ""
        for ext in extensions + [None]:  # Try last with no extension
            for converter in self._page_converters:
                _kwargs = copy.deepcopy(kwargs)

                # Overwrite file_extension appropriately
                if ext is None:
                    if "file_extension" in _kwargs:
                        del _kwargs["file_extension"]
                else:
                    _kwargs.update({"file_extension": ext})

                # Copy any additional global options
                if "mlm_client" not in _kwargs and self._mlm_client is not None:
                    _kwargs["mlm_client"] = self._mlm_client

                if "mlm_model" not in _kwargs and self._mlm_model is not None:
                    _kwargs["mlm_model"] = self._mlm_model

                # If we hit an error log it and keep trying
                try:
                    res = converter.convert(local_path, **_kwargs)
                except Exception:
                    error_trace = ("\n\n" + traceback.format_exc()).strip()

                if res is not None:
                    # Normalize the content
                    res.text_content = "\n".join([line.rstrip() for line in re.split(r"\r?\n", res.text_content)])
                    res.text_content = re.sub(r"\n{3,}", "\n\n", res.text_content)

                    # Todo
                    return res

        # If we got this far without success, report any exceptions
        if len(error_trace) > 0:
            raise FileConversionException(
                f"Could not convert '{local_path}' to Markdown. File type was recognized as {extensions}. While converting the file, the following error was encountered:\n\n{error_trace}"
            )

        # Nothing can handle it!
        raise UnsupportedFormatException(
            f"Could not convert '{local_path}' to Markdown. The formats {extensions} are not supported."
        )

    def _append_ext(self, extensions, ext):
        """Append a unique non-None, non-empty extension to a list of extensions."""
        if ext is None:
            return
        ext = ext.strip()
        if ext == "":
            return
        # if ext not in extensions:
        if True:
            extensions.append(ext)

    def _guess_ext_magic(self, path):
        """Use puremagic (a Python implementation of libmagic) to guess a file's extension based on the first few bytes."""
        # Use puremagic to guess
        try:
            guesses = puremagic.magic_file(path)
            if len(guesses) > 0:
                ext = guesses[0].extension.strip()
                if len(ext) > 0:
                    return ext
        except FileNotFoundError:
            pass
        except IsADirectoryError:
            pass
        except PermissionError:
            pass
        return None

    def register_page_converter(self, converter: DocumentConverter) -> None:
        """Register a page text converter."""
        self._page_converters.insert(0, converter)
